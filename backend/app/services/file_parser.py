from dataclasses import dataclass
from pathlib import Path

import fitz
import openpyxl
from docx import Document
from pptx import Presentation

from app.config import DATA_DIR, settings


@dataclass
class ParseResult:
    parsed_text_path: str | None
    status: str
    original_chars: int = 0
    used_chars: int = 0
    is_truncated: bool = False


def parse_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding, errors="ignore")
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="ignore")


def parse_pdf(path: Path) -> str:
    parts: list[str] = []
    with fitz.open(path) as doc:
        for index, page in enumerate(doc, start=1):
            parts.append(f"\n--- 第 {index} 页 ---\n{page.get_text()}")
    return "\n".join(parts)


def parse_docx(path: Path) -> str:
    doc = Document(path)
    lines = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(lines)


def parse_pptx(path: Path) -> str:
    prs = Presentation(path)
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n--- 幻灯片 {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def parse_xlsx(path: Path) -> str:
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"\n--- 工作表：{sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                lines.append("\t".join(values))
    return "\n".join(lines)


def parse_file(path: str, attachment_id: str) -> ParseResult:
    source = Path(path)
    ext = source.suffix.lower()
    parsed_dir = DATA_DIR / "parsed"
    parsed_dir.mkdir(parents=True, exist_ok=True)

    if ext in {".txt", ".md", ".csv"}:
        text = parse_text_file(source)
    elif ext == ".pdf":
        text = parse_pdf(source)
    elif ext == ".docx":
        text = parse_docx(source)
    elif ext == ".pptx":
        text = parse_pptx(source)
    elif ext == ".xlsx":
        text = parse_xlsx(source)
    else:
        return ParseResult(None, "unsupported")

    used_text = text[: settings.max_file_chars]
    parsed_path = parsed_dir / f"{attachment_id}.txt"
    parsed_path.write_text(used_text, encoding="utf-8")
    return ParseResult(
        parsed_text_path=str(parsed_path),
        status="parsed",
        original_chars=len(text),
        used_chars=len(used_text),
        is_truncated=len(text) > len(used_text),
    )


def read_parsed_text(path: str | None) -> str:
    if not path:
        return ""
    parsed = Path(path)
    if not parsed.exists():
        return ""
    return parsed.read_text(encoding="utf-8", errors="ignore")
